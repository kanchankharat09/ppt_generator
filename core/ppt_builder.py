from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from core.image_utils import get_image_size_inches
from core.schemas import Slide, SlidePlan

CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}

SLIDE_WIDTH_IN = 10
SLIDE_HEIGHT_IN = 7.5


def _add_images_to_right_column(slide, image_filenames, image_store, top_in=2.0):
    """Places one or more images stacked in the right half of the slide, each scaled
    to fit its box without stretching, so it never overlaps the left-column content."""
    box_width_in = 4.0
    box_height_in = (SLIDE_HEIGHT_IN - top_in - 0.5) / max(len(image_filenames), 1)
    left_in = SLIDE_WIDTH_IN - box_width_in - 0.5

    current_top = top_in
    for filename in image_filenames:
        image_bytes = image_store.get(filename)
        if image_bytes is None:
            continue

        width_in, height_in = get_image_size_inches(image_bytes, box_width_in, box_height_in)
        slide.shapes.add_picture(
            BytesIO(image_bytes), Inches(left_in), Inches(current_top),
            width=Inches(width_in), height=Inches(height_in),
        )
        current_top += box_height_in


def _add_bullet_slide(prs: Presentation, bullet_layout, title_only_layout, slide_data: Slide, image_store: dict):
    has_images = bool(slide_data.image_filenames)

    if not has_images:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_data.title

        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()

        for i, bullet in enumerate(slide_data.bullets):
            if i == 0:
                text_frame.text = bullet
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = bullet
        return

    # With images, use a title-only layout and lay out bullets (left) + images (right)
    # ourselves, so the image never overlaps the text placeholder.
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = slide_data.title

    if slide_data.bullets:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5.3))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        for i, bullet in enumerate(slide_data.bullets):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            paragraph.text = f"• {bullet}"
            paragraph.font.size = Pt(16)

    _add_images_to_right_column(slide, slide_data.image_filenames, image_store, top_in=1.5)


def _add_chart_slide(prs: Presentation, title_only_layout, slide_data: Slide, image_store: dict):
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = slide_data.title

    has_images = bool(slide_data.image_filenames)
    chart_width_in = 5.5 if has_images else 8
    chart_left_in = 0.5

    # Bullets go in a small text box above the chart, since this layout has no body placeholder.
    if slide_data.bullets:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(chart_width_in + 0.5), Inches(1))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        for i, bullet in enumerate(slide_data.bullets):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            paragraph.text = f"• {bullet}"
            paragraph.font.size = Pt(14)

    chart = slide_data.chart
    chart_data = CategoryChartData()
    chart_data.categories = chart.categories
    chart_data.add_series(chart.series_name, chart.values)

    xl_chart_type = CHART_TYPE_MAP.get(chart.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    slide.shapes.add_chart(
        xl_chart_type, Inches(chart_left_in), Inches(2.4),
        Inches(chart_width_in), Inches(4.2), chart_data,
    )

    if has_images:
        _add_images_to_right_column(slide, slide_data.image_filenames, image_store, top_in=2.4)


def build_pptx(plan: SlidePlan, image_store: dict | None = None) -> BytesIO:
    image_store = image_store or {}
    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = plan.presentation_title
    slide.placeholders[1].text = "Generated with AI PPT Generator"

    if plan.title_slide_image_filenames:
        _add_images_to_right_column(slide, plan.title_slide_image_filenames, image_store, top_in=3.0)

    bullet_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]

    for slide_data in plan.slides:
        if slide_data.chart is not None:
            _add_chart_slide(prs, title_only_layout, slide_data, image_store)
        else:
            _add_bullet_slide(prs, bullet_layout, title_only_layout, slide_data, image_store)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
