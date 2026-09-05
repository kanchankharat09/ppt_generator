from pptx.dml.color import RGBColor

THEMES = {
    "Default": {
        "background": None,  # keeps PowerPoint's plain white default
        "title_color": RGBColor(0x00, 0x00, 0x00),
        "body_color": RGBColor(0x33, 0x33, 0x33),
        "font_name": "Calibri",
    },
    "Professional Blue": {
        "background": RGBColor(0xF3, 0xF7, 0xFB),
        "title_color": RGBColor(0x1B, 0x3A, 0x5C),
        "body_color": RGBColor(0x2E, 0x2E, 0x2E),
        "font_name": "Calibri",
    },
    "Minimal Dark": {
        "background": RGBColor(0x1E, 0x1E, 0x1E),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
        "body_color": RGBColor(0xE0, 0xE0, 0xE0),
        "font_name": "Calibri",
    },
    "Warm Neutral": {
        "background": RGBColor(0xFA, 0xF3, 0xE9),
        "title_color": RGBColor(0x5C, 0x3A, 0x1B),
        "body_color": RGBColor(0x3D, 0x3D, 0x3D),
        "font_name": "Georgia",
    },
}


def apply_background(slide, theme: dict):
    if theme["background"] is not None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = theme["background"]


def style_title(slide, theme: dict):
    if slide.shapes.title is None:
        return
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        paragraph.font.color.rgb = theme["title_color"]
        paragraph.font.name = theme["font_name"]


def style_text_frame(text_frame, theme: dict):
    for paragraph in text_frame.paragraphs:
        paragraph.font.color.rgb = theme["body_color"]
        paragraph.font.name = theme["font_name"]
